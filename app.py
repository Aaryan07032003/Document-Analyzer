import os
import json
from flask import Flask, render_template, request, jsonify
import google.generativeai as genai
from PyPDF2 import PdfReader
from PIL import Image
import pytesseract
from pptx import Presentation
import io
from werkzeug.utils import secure_filename

app = Flask(__name__)

UPLOAD_FOLDER = 'temp_uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

class DocumentAnalyzer:
    def __init__(self):
        self.documents = []
        self.analyzed_content = ""
        self.knowledge_base = self.load_knowledge_base()
        self.init_gemini()

    def init_gemini(self):
        genai.configure(api_key="AIzaSyCg2Bxk6EMDcRgd_1CClj8jhDhkS831LVk")
        self.model = genai.GenerativeModel('gemini-pro')

    def load_knowledge_base(self):
        try:
            with open('knowledge_base.json', 'r') as f:
                return json.load(f)
        except FileNotFoundError:
            return {}

    def save_knowledge_base(self):
        with open('knowledge_base.json', 'w') as f:
            json.dump(self.knowledge_base, f)

    def process_document(self, file, filename):
        file_extension = os.path.splitext(filename)[1].lower()
        try:
            if file_extension == '.pdf':
                return self.process_pdf(file, filename)
            elif file_extension in ['.png', '.jpg', '.jpeg']:
                return self.process_image(file, filename)
            elif file_extension in ['.ppt', '.pptx']:
                return self.process_ppt(file, filename)
            else:
                content = file.read().decode('utf-8')
                print(f"Processed {filename}: {len(content)} characters") 
                return content
        except Exception as e:
            print(f"Error processing {filename}: {str(e)}")  
            return f"Error processing {filename}: {str(e)}"

    def process_pdf(self, file, filename):
        try:
            pdf = PdfReader(file)
            text = ""
            for page in pdf.pages:
                text += page.extract_text() + "\n"
            print(f"Extracted {len(text)} characters from PDF {filename}")  
            return text
        except Exception as e:
            print(f"Error processing PDF {filename}: {str(e)}")  
            return ""

    def process_image(self, file, filename):
        try:
            text = pytesseract.image_to_string(Image.open(file))
            print(f"Extracted {len(text)} characters from image {filename}")  
            return text
        except Exception as e:
            print(f"Error processing image {filename}: {str(e)}")  
            return ""

    def process_ppt(self, file, filename):
        try:
            prs = Presentation(file)
            text = " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text'))
            print(f"Extracted {len(text)} characters from PPT {filename}")  
            return text
        except Exception as e:
            print(f"Error processing PPT {filename}: {str(e)}")  
            return ""

    def analyze_documents(self):
        if not self.documents:
            return "Please upload documents first."

        all_text = ""
        for doc_path in self.documents:
            with open(doc_path, 'rb') as file:
                filename = os.path.basename(doc_path)
                processed_text = self.process_document(file, filename)
                print(f"Processed {filename}: {len(processed_text)} characters") 
                all_text += processed_text + "\n\n"

        self.analyzed_content = all_text
        print(f"Total processed text: {len(all_text)} characters")  

        analysis_prompt = """
        Perform a comprehensive analysis of the following documents. Your analysis should be structured in numbered points, with each point on a new line. Include:

        1. Key points and main themes
        2. Important details and insightful interpretations
        3. Potential trends or patterns
        4. Significant data or statistics with context
        5. Connections between different pieces of information
        6. Possible implications or future scenarios based on the content
        7. Gaps in the information or areas that require further investigation

        Approach this task with a deep, analytical mindset. Go beyond surface-level observations to provide a nuanced understanding of the content. Format your response as a numbered list, with each main point and its sub-points on separate lines.

        Document Contents:
        """

        full_prompt = f"{analysis_prompt}\n{all_text}"

        try:
            response = self.model.generate_content(full_prompt)
            analysis_result = response.text
            
            formatted_result = self.format_analysis_result(analysis_result)
            
            self.knowledge_base[','.join([os.path.basename(doc) for doc in self.documents])] = formatted_result
            self.save_knowledge_base()
            
            return formatted_result
        except Exception as e:
            print(f"Error in analysis: {str(e)}")  
            return f"Error in analysis: {str(e)}"

    def format_analysis_result(self, result):
        lines = result.split('\n')
        formatted_lines = []
        for line in lines:
            if line.strip() and line.strip()[0].isdigit() and '.' in line:
                formatted_lines.append('\n' + line.strip())
            else:
                formatted_lines.append(line.strip())
        return '\n'.join(formatted_lines).strip()

    def get_answer(self, question):
        if not self.analyzed_content:
            return "Please analyze documents first."

        if not question:
            return "Please enter a question."

        print(f"Analyzed content length: {len(self.analyzed_content)}")  

        qa_prompt = f"""
        Based on the following analyzed content and additional knowledge, please answer the question. Your response should:

        1. Demonstrate a deep understanding of the content and its implications.
        2. Draw from multiple pieces of information to form a comprehensive answer.
        3. If the question is about future predictions, trends, or potential outcomes:
           - Use the provided information, historical data, and logical reasoning to make informed projections.
           - Consider multiple possible scenarios and their likelihood.
           - Explain the factors that could influence these future outcomes.
           - Provide a timeframe for your predictions when possible.
           - If exact future data isn't available, provide educated estimates based on historical trends and current data.
           - Clearly state when you're making an estimate and explain your reasoning.
        4. Consider potential biases or limitations in the data when making predictions or drawing conclusions.
        5. Provide reasoning and justification for your answers and predictions.
        6. If the question is not explicitly about the future, still consider relevant future implications if applicable.

        Remember, it's okay to make reasonable estimates based on the data provided. When doing so, explain your reasoning and any assumptions made.

        Analyzed Content:
        {self.analyzed_content}

        Additional Knowledge:
        {json.dumps(self.knowledge_base)}

        Question: {question}

        Answer:
        """

        try:
            response = self.model.generate_content(qa_prompt)
            formatted_result = self.format_analysis_result(response.text)
            return formatted_result
        except Exception as e:
            print(f"Error in generating answer: {str(e)}")  
            return f"Error in generating answer: {str(e)}"

analyzer = DocumentAnalyzer()

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload():
    if 'files[]' not in request.files:
        return jsonify({"error": "No file part"})
    
    files = request.files.getlist('files[]')
    
    if not files or files[0].filename == '':
        return jsonify({"error": "No selected file"})
    
    uploaded_files = []
    for file in files:
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        uploaded_files.append(file_path)
    
    analyzer.documents = uploaded_files
    print(f"Uploaded {len(uploaded_files)} files: {uploaded_files}")  
    return jsonify({"message": f"{len(uploaded_files)} documents uploaded successfully"})

@app.route('/analyze', methods=['POST'])
def analyze():
    try:
        result = analyzer.analyze_documents()
        return jsonify({"result": result})
    except Exception as e:
        print(f"Error in analysis: {str(e)}")  
        return jsonify({"error": f"An error occurred during analysis: {str(e)}"})

@app.route('/answer', methods=['POST'])
def answer():
    question = request.json.get('question')
    try:
        result = analyzer.get_answer(question)
        return jsonify({"result": result})
    except Exception as e:
        print(f"Error in generating answer: {str(e)}")  
        return jsonify({"error": f"An error occurred while generating the answer: {str(e)}"})

if __name__ == '__main__':
    app.run(debug=True)
